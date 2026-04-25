#!/usr/bin/env python3
"""doc2md — PDF / DOCX / PPTX → Markdown converter (compatible with md2pdf)."""

from __future__ import annotations

import argparse
import os
import re
import sys
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional, Tuple

# ---------------------------------------------------------------------------
# DocToken — universal intermediate token
# ---------------------------------------------------------------------------

@dataclass
class DocToken:
    kind: str  # heading, paragraph, table, image, code, list, chart, quote, hr, footnote
    text: str = ""
    level: int = 0
    rows: list = field(default_factory=list)
    headers: list = field(default_factory=list)
    aligns: list = field(default_factory=list)
    caption: str = ""
    image_path: str = ""
    image_data: bytes = b""
    lang: str = ""
    items: list = field(default_factory=list)  # for list: [(level, text, ordered)]
    chart_spec: str = ""  # raw ```chart block content


# ---------------------------------------------------------------------------
# Common utilities
# ---------------------------------------------------------------------------

def _merge_lists(tokens: List[DocToken]) -> List[DocToken]:
    """Merges consecutive list tokens with improved nesting handling (#3).

    Each item carries its own (level, text, ordered) — merge if:
    - sub-item (level > 0),
    - previous list contains sub-items,
    - or same ordered type at the base level.
    """
    merged: List[DocToken] = []
    for tok in tokens:
        if tok.kind == "list" and merged and merged[-1].kind == "list":
            cur_items = tok.items
            prev_items = merged[-1].items
            if cur_items and prev_items:
                cur_level, _, cur_ordered = cur_items[0]
                # Base type of previous list (first item with level==0)
                base_ordered = next(
                    (o for l, _, o in prev_items if l == 0),
                    prev_items[0][2],
                )
                # Merge if: sub-item, or returning from sub-items,
                # or same ordered type at the base level
                prev_last_level = prev_items[-1][0]
                if cur_level > 0 or prev_last_level > 0 or cur_ordered == base_ordered:
                    merged[-1].items.extend(tok.items)
                    continue
        merged.append(tok)
    return merged


# ---------------------------------------------------------------------------
# Code language detection (#7)
# ---------------------------------------------------------------------------

def detect_code_language(code: str) -> str:
    """Heuristic detection of code language by content."""
    if not code or len(code.strip()) < 10:
        return ""

    text = code.strip()

    # Python
    py = 0
    if re.search(r'^(def |class |import |from \S+ import )', text, re.M):
        py += 3
    if re.search(r'(print\(|self\.|__\w+__|\.append\(|\.items\(\))', text):
        py += 2
    if re.search(r'^(if .+:|elif |else:|for .+ in |while .+:|try:|except |with .+ as )', text, re.M):
        py += 1
    if re.search(r'(\.py\b|#!.*python|# -\*- coding)', text):
        py += 2

    # JavaScript / TypeScript
    js = 0
    if re.search(r'\b(const |let |var |function |=>|require\(|module\.exports)', text):
        js += 3
    if re.search(r'(console\.(log|error|warn)|document\.|window\.|\.addEventListener)', text):
        js += 2
    if re.search(r'\b(async |await |Promise|\.then\(|\.catch\()', text):
        js += 1

    # Java
    java = 0
    if re.search(r'\b(public |private |protected |class |interface |void |static )', text):
        java += 3
    if re.search(r'(System\.out\.print|\.println\(|@Override|@Nullable)', text):
        java += 2
    if re.search(r'\bnew \w+\(', text):
        java += 1

    # C / C++
    c = 0
    if re.search(r'^#include\s*[<"]', text, re.M):
        c += 4
    if re.search(r'\b(int main\(|printf\(|sizeof\(|malloc\(|free\()', text):
        c += 2
    if re.search(r'(std::|cout|cin|endl|nullptr|#define )', text):
        c += 2

    # SQL
    sql = 0
    if re.search(r'\b(SELECT|INSERT|UPDATE|DELETE|CREATE TABLE|ALTER TABLE|DROP)\b', text):
        sql += 3
    if re.search(r'\b(FROM|WHERE|JOIN|GROUP BY|ORDER BY|HAVING|LIMIT)\b', text):
        sql += 2

    # Shell / Bash
    sh = 0
    if re.search(r'^(#!/bin/(ba)?sh|#!/usr/bin/env bash)', text, re.M):
        sh += 4
    if re.search(r'(\$\{?\w+\}?|\becho\b|\bif \[|\bfi\b|\bdone\b|\bdo\b)', text):
        sh += 2
    if re.search(r'(\|\s*grep|\|\s*awk|\|\s*sed|&&\s*\w)', text):
        sh += 1

    # HTML / XML
    html = 0
    if re.search(r'<(html|div|span|body|head|p|table|form)\b', text, re.I):
        html += 3
    if re.search(r'<\w+[^>]*>.*</\w+>', text, re.S):
        html += 1

    # CSS
    css = 0
    if re.search(r'[.#]\w+\s*\{[^}]*(color|margin|padding|font|display|background)', text, re.S):
        css += 3
    if re.search(r'@media|@import|@keyframes', text):
        css += 2

    # JSON
    json_s = 0
    stripped = text.strip()
    if (stripped.startswith("{") and stripped.endswith("}")) or \
       (stripped.startswith("[") and stripped.endswith("]")):
        if re.search(r'"\w+":\s*["\d\[{tfn]', text):
            json_s += 4

    # YAML
    yaml_s = 0
    if re.search(r'^\w[\w-]*:\s+\S', text, re.M) and not re.search(r'[{};]', text):
        yaml_s += 2
    if re.search(r'^---\s*$', text, re.M):
        yaml_s += 1

    # Swift (more specific patterns to avoid conflicting with Go)
    swift = 0
    if re.search(r'\b(struct \w+|enum \w+|protocol \w+|guard let|@IBOutlet|@IBAction|@objc)', text):
        swift += 3
    if re.search(r'(\.self|\.Type|override func|import (UIKit|Foundation|SwiftUI))', text):
        swift += 2
    if re.search(r'\b(var \w+:\s|let \w+:\s|func \w+\(.*\)\s*->\s)', text):
        swift += 1

    # Go
    go = 0
    if re.search(r'^package \w+', text, re.M):
        go += 3
    if re.search(r'\bfunc \w+\(', text):
        go += 2
    if re.search(r'(:= |fmt\.|go func|make\(|range |defer )', text):
        go += 2

    # Rust
    rust = 0
    if re.search(r'\b(fn \w+|let mut |impl |pub fn |use std::)', text):
        rust += 3
    if re.search(r'(-> |&str|&self|println!|vec!|Option<|Result<)', text):
        rust += 2

    # R
    r_s = 0
    if re.search(r'(<-\s|library\(|ggplot\(|data\.frame\(|summary\()', text):
        r_s += 3

    scores = {
        "python": py, "javascript": js, "java": java, "c": c,
        "sql": sql, "bash": sh, "html": html, "css": css,
        "json": json_s, "yaml": yaml_s, "swift": swift,
        "go": go, "rust": rust, "r": r_s,
    }
    best = max(scores, key=scores.get)
    if scores[best] >= 3:
        return best
    return ""


# ---------------------------------------------------------------------------
# PDFConverter (#1, #2, #9, #10)
# ---------------------------------------------------------------------------

class PDFConverter:
    """Converts PDF → list of DocToken."""

    # Keywords for monospace fonts (used to detect code)
    MONO_KEYWORDS = (
        "courier", "mono", "consolas", "menlo", "source code",
        "fira code", "inconsolata", "dejavu sans mono", "liberation mono",
        "roboto mono", "jetbrains", "hack", "iosevka",
    )

    def convert(self, pdf_path: str, image_dir: str, extract_images: bool = True,
                verbose: bool = False) -> Tuple[List[DocToken], dict]:
        try:
            import fitz
        except ImportError:
            print("ERROR: PyMuPDF is not installed. pip install PyMuPDF", file=sys.stderr)
            sys.exit(1)

        doc = fitz.open(pdf_path)
        tokens: List[DocToken] = []
        meta = {}

        # --- Extended metadata (#8) ---
        m = doc.metadata or {}
        if m.get("title") and m["title"] not in ("(anonymous)", "untitled", ""):
            meta["title"] = m["title"]
        if m.get("author") and m["author"] not in ("(anonymous)", ""):
            meta["author"] = m["author"]
        if m.get("subject"):
            meta["subject"] = m["subject"]
        if m.get("keywords"):
            meta["keywords"] = m["keywords"]
        if m.get("creationDate"):
            d = self._parse_pdf_date(m["creationDate"])
            if d:
                meta["date"] = d
        if m.get("creator"):
            meta["creator"] = m["creator"]

        # ═══════════════════════════════════════════════════════════════
        # Single pass: cache all data (#10)
        # ═══════════════════════════════════════════════════════════════
        all_sizes: List[float] = []
        font_counter: Counter = Counter()
        page_cache = []  # [(page_num, cached_blocks, table_tokens, image_tokens)]
        page_widths: List[float] = []
        img_counter = 0

        for page_num, page in enumerate(doc):
            pw = page.rect.width
            page_widths.append(pw)

            # --- Tables ---
            table_rects = []
            table_tokens = []
            try:
                tables = page.find_tables()
                for table in tables:
                    table_rects.append(fitz.Rect(table.bbox))
                    rows_data = []
                    for row in table.extract():
                        rows_data.append([cell if cell else "" for cell in row])
                    if rows_data:
                        hdrs = rows_data[0]
                        body_rows = rows_data[1:] if len(rows_data) > 1 else []
                        table_tokens.append(DocToken(kind="table", headers=hdrs, rows=body_rows))
            except Exception:
                pass

            # --- Images ---
            image_tokens = []
            if extract_images:
                for img_info in page.get_images(full=True):
                    xref = img_info[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n > 4:
                            pix = fitz.Pixmap(fitz.csRGB, pix)
                        if pix.width >= 50 and pix.height >= 50:
                            img_counter += 1
                            ext = "png"
                            img_name = f"img_{page_num + 1}_{img_counter}.{ext}"
                            img_path = os.path.join(image_dir, img_name)
                            pix.save(img_path)
                            image_tokens.append(DocToken(kind="image", image_path=img_name))
                        pix = None
                    except Exception:
                        pass

            # --- Text blocks: cache raw data ---
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
            cached_blocks = []
            for b in blocks:
                if b["type"] != 0:
                    continue
                block_rect = fitz.Rect(b["bbox"])
                if any(block_rect.intersects(tr) for tr in table_rects):
                    continue

                lines_data = []
                for line in b["lines"]:
                    spans = []
                    for span in line["spans"]:
                        sd = {
                            "text": span["text"],
                            "size": round(span["size"], 1),
                            "flags": span["flags"],
                            "font": span.get("font", ""),
                        }
                        txt = span["text"].strip()
                        if txt:
                            all_sizes.append(sd["size"])
                            font_counter[span.get("font", "")] += len(txt)
                        spans.append(sd)
                    lines_data.append(spans)

                cached_blocks.append({
                    "bbox": b["bbox"],
                    "lines": lines_data,
                })

            page_cache.append((page_num, cached_blocks, table_tokens, image_tokens))

        if not all_sizes:
            doc.close()
            return tokens, meta

        # ═══════════════════════════════════════════════════════════════
        # Compute metrics from cache (#10)
        # ═══════════════════════════════════════════════════════════════
        body_size = max(set(all_sizes), key=all_sizes.count)

        # Detect monospace fonts
        mono_fonts = set()
        for fname in font_counter:
            if any(kw in fname.lower() for kw in self.MONO_KEYWORDS):
                mono_fonts.add(fname)

        avg_page_width = sum(page_widths) / len(page_widths) if page_widths else 612.0

        # ═══════════════════════════════════════════════════════════════
        # Process cached blocks
        # ═══════════════════════════════════════════════════════════════
        for page_num, cached_blocks, table_tokens, image_tokens in page_cache:
            tokens.extend(table_tokens)
            tokens.extend(image_tokens)

            # Column detection (#2) — per-page
            ordered_blocks = cached_blocks
            if len(cached_blocks) >= 4:
                ordered_blocks = self._sort_by_columns(cached_blocks, avg_page_width)

            for block in ordered_blocks:
                self._process_cached_block(block, tokens, body_size, mono_fonts)

        doc.close()

        # Post-processing
        tokens = self._filter_noise(tokens)
        tokens = _merge_lists(tokens)
        tokens = self._merge_code_blocks(tokens)
        return tokens, meta

    # --- Multi-factor heading detection (#1) ---
    @staticmethod
    def _heading_score(block_size: float, body_size: float,
                       is_bold: bool, text: str) -> Tuple[int, float]:
        """Returns (heading_level, confidence). Level 0 = not a heading."""
        if not text.strip():
            return 0, 0.0

        score = 0.0
        clean = re.sub(r'\*+', '', text).strip()

        # Factor 1: font size (most significant)
        size_ratio = block_size / body_size if body_size > 0 else 1.0
        if size_ratio >= 1.8:
            score += 5.0
        elif size_ratio >= 1.4:
            score += 4.0
        elif size_ratio >= 1.2:
            score += 3.0
        elif size_ratio >= 1.08:
            score += 2.0

        # Factor 2: bold font
        if is_bold:
            score += 1.5

        # Factor 3: short text (headings are typically < 120 chars)
        if len(clean) < 80:
            score += 1.0
        elif len(clean) < 120:
            score += 0.5

        # Factor 4: ALL UPPERCASE or Title Case
        if clean.isupper() and len(clean) > 2:
            score += 1.0
        elif clean.istitle() and len(clean.split()) <= 8:
            score += 0.5

        # Factor 5: no terminal sentence punctuation
        if not re.search(r'[.;,]$', clean):
            score += 0.5

        # Anti-factor: TOC fillers (dots, ellipses) — not headings
        dots_ratio = sum(1 for c in clean if c in '·…•.') / max(len(clean), 1)
        if dots_ratio > 0.3:
            return 0, 0.0

        # Factor 6: single line
        if '\n' not in text and len(clean) < 200:
            score += 0.3

        # Determine level
        if score >= 6.0 and size_ratio >= 1.5:
            return 1, score
        elif score >= 4.5 and size_ratio >= 1.2:
            return 2, score
        elif score >= 3.5:
            return 3, score
        elif score >= 2.5 and (is_bold or size_ratio >= 1.08):
            return 4, score

        return 0, score

    def _process_cached_block(self, block: dict, tokens: list,
                              body_size: float, mono_fonts: set):
        """Processes a cached text block → DocToken(s)."""
        lines_data = block["lines"]

        # Collect per-span info
        block_parts = []
        span_sizes: List[float] = []
        span_bolds: List[bool] = []
        all_mono = True
        any_mono = False
        total_chars = 0

        for line_spans in lines_data:
            line_parts = []
            for sd in line_spans:
                txt = sd["text"]
                if not txt.strip():
                    line_parts.append(txt)
                    continue

                sz = sd["size"]
                flags = sd["flags"]
                font = sd["font"]
                is_bold = bool(flags & (1 << 4))
                is_italic = bool(flags & (1 << 1))
                is_super = bool(flags & (1 << 0))
                is_mono = font in mono_fonts

                span_sizes.append(sz)
                span_bolds.append(is_bold)
                total_chars += len(txt.strip())

                if is_mono:
                    any_mono = True
                else:
                    all_mono = False

                # Inline formatting (#9)
                t = txt
                if is_mono and not all_mono:
                    # Inline code inside a paragraph
                    t = f"`{t.strip()}`"
                elif is_bold and is_italic:
                    t = f"***{t.strip()}***"
                elif is_bold:
                    t = f"**{t.strip()}**"
                elif is_italic:
                    t = f"*{t.strip()}*"
                if is_super:
                    t = f"^{t.strip()}^"

                line_parts.append(t)

            line_text = "".join(line_parts).rstrip()
            if line_text:
                block_parts.append(line_text)

        if not block_parts:
            return

        full_text = " ".join(block_parts).strip()
        if not full_text:
            return

        # Code block: all spans are monospace
        if all_mono and any_mono and total_chars > 5:
            raw_lines = []
            for line_spans in lines_data:
                parts = [sd["text"] for sd in line_spans]
                raw_lines.append("".join(parts).rstrip())
            code_text = "\n".join(raw_lines).strip()
            if code_text:
                tokens.append(DocToken(kind="code", text=code_text))
                return

        # List detection
        list_match = re.match(r'^([•●○▪▸\-–—])\s+(.+)', full_text)
        num_match = re.match(r'^(\d+)[.)]\s+(.+)', full_text)

        if list_match:
            tokens.append(DocToken(kind="list", items=[(0, list_match.group(2), False)]))
        elif num_match:
            tokens.append(DocToken(kind="list", items=[(0, num_match.group(2), True)]))
        else:
            # Multi-factor heading detection (#1)
            if span_sizes:
                size_counter = Counter(span_sizes)
                block_size = size_counter.most_common(1)[0][0]
            else:
                block_size = body_size

            is_block_bold = (
                sum(1 for b in span_bolds if b) > len(span_bolds) / 2
                if span_bolds else False
            )

            level, _score = self._heading_score(
                block_size, body_size, is_block_bold, full_text,
            )

            if level > 0:
                clean = re.sub(r'\*+', '', full_text).strip()
                tokens.append(DocToken(kind="heading", text=clean, level=level))
            else:
                tokens.append(DocToken(kind="paragraph", text=full_text))

    # --- Column detection (#2) ---
    @staticmethod
    def _sort_by_columns(blocks: list, page_width: float) -> list:
        """Detects multi-column layout and sorts blocks by columns."""
        if not blocks:
            return blocks

        x_positions = [round(b["bbox"][0]) for b in blocks]
        if not x_positions:
            return blocks

        margin = page_width * 0.1
        inner_x = sorted(set(x for x in x_positions if x > margin * 0.5))

        if len(inner_x) < 2:
            return blocks

        # Clustering: group x-positions within 20pt
        clusters = []
        current_cluster = [inner_x[0]]
        for x in inner_x[1:]:
            if x - current_cluster[-1] <= 20:
                current_cluster.append(x)
            else:
                clusters.append(sum(current_cluster) / len(current_cluster))
                current_cluster = [x]
        clusters.append(sum(current_cluster) / len(current_cluster))

        if len(clusters) < 2:
            return blocks

        # Only handle 2-column layout
        if len(clusters) == 2:
            gap = clusters[1] - clusters[0]
            if gap < page_width * 0.2 or gap > page_width * 0.7:
                return blocks  # too close or too far — not columns

            mid = (clusters[0] + clusters[1]) / 2.0

            left_blocks = []
            right_blocks = []
            for b in blocks:
                center_x = (b["bbox"][0] + b["bbox"][2]) / 2.0
                if center_x < mid:
                    left_blocks.append(b)
                else:
                    right_blocks.append(b)

            left_blocks.sort(key=lambda b: b["bbox"][1])
            right_blocks.sort(key=lambda b: b["bbox"][1])
            return left_blocks + right_blocks

        # 3+ clusters — sort by x then by y
        return sorted(blocks, key=lambda b: (b["bbox"][0] // 100, b["bbox"][1]))

    # --- Merging code blocks ---
    @staticmethod
    def _merge_code_blocks(tokens: List[DocToken]) -> List[DocToken]:
        """Merges consecutive code blocks."""
        merged = []
        for tok in tokens:
            if tok.kind == "code" and merged and merged[-1].kind == "code":
                merged[-1].text += "\n" + tok.text
                continue
            merged.append(tok)
        return merged

    @staticmethod
    def _parse_pdf_date(date_str: str) -> str:
        """Parses PDF date (D:YYYYMMDDHHmmSS) → ISO."""
        if not date_str:
            return ""
        d = date_str.replace("D:", "").strip()
        try:
            if len(d) >= 8:
                return f"{d[0:4]}-{d[4:6]}-{d[6:8]}"
        except Exception:
            pass
        return ""

    @staticmethod
    def _filter_noise(tokens: List[DocToken]) -> List[DocToken]:
        """Filters out watermarks, page numbers, repeating headers/footers."""
        if not tokens:
            return tokens

        text_freq: dict = {}
        for tok in tokens:
            if tok.kind == "paragraph" and tok.text.strip():
                key = re.sub(r'\*+', '', tok.text).strip()
                text_freq[key] = text_freq.get(key, 0) + 1

        noise_texts = {text for text, count in text_freq.items() if count >= 3}

        filtered = []
        for tok in tokens:
            if tok.kind == "paragraph":
                clean = re.sub(r'\*+', '', tok.text).strip()
                if clean in noise_texts:
                    continue
                if re.match(r'^\d{1,4}$', clean):
                    continue
            filtered.append(tok)

        return filtered


# ---------------------------------------------------------------------------
# HtmlDocConverter — for HTML files disguised as .docx
# ---------------------------------------------------------------------------

class HtmlDocConverter:
    """Converts HTML (disguised as .docx) → list of DocToken."""

    def convert(self, html_path: str, image_dir: str, extract_images: bool = True,
                verbose: bool = False) -> Tuple[List[DocToken], dict]:
        from lxml import html as lxml_html

        with open(html_path, "r", encoding="utf-8") as f:
            content = f.read()

        tree = lxml_html.fromstring(content)
        tokens: List[DocToken] = []
        meta = {}

        # Extract title from <title> tag
        title_el = tree.find(".//title")
        if title_el is not None and title_el.text:
            raw_title = title_el.text.strip()
            for ext in (".docx", ".doc", ".html", ".htm"):
                if raw_title.lower().endswith(ext):
                    raw_title = raw_title[:-len(ext)].strip()
            if raw_title and not re.match(r'^\d+$', raw_title) and len(raw_title) > 3:
                meta["title"] = raw_title

        class_props = self._parse_css_classes(content)

        body = tree.find(".//body")
        if body is None:
            return tokens, meta

        size_counts = {}
        for span in body.iter("span"):
            cls = span.get("class", "")
            txt = (span.text or "").strip()
            if cls in class_props and class_props[cls]["font_size"] and txt:
                sz = class_props[cls]["font_size"]
                size_counts[sz] = size_counts.get(sz, 0) + len(txt)

        if size_counts:
            self._body_size = max(size_counts, key=size_counts.get)
        else:
            self._body_size = 16
        self._h1_thresh = self._body_size * 1.8
        self._h2_thresh = self._body_size * 1.4

        self._process_element(body, tokens, class_props, image_dir, extract_images)

        tokens = _merge_lists(tokens)
        return tokens, meta

    def _parse_css_classes(self, html_content: str) -> dict:
        """Parse CSS class definitions → {class_name: {font_size, bold, italic, ...}}."""
        props = {}
        for m in re.finditer(r'\.(\w+)\s*\{([^}]+)\}', html_content):
            cls_name = m.group(1)
            css_text = m.group(2)

            info = {"bold": False, "italic": False, "font_size": None, "text_align": None}

            if "font-weight:bold" in css_text or "font-weight: bold" in css_text:
                info["bold"] = True
            if "font-style:italic" in css_text or "font-style: italic" in css_text:
                info["italic"] = True

            sz_match = re.search(r'font-size:\s*([\d.]+)(em|px|pt|;)', css_text)
            if sz_match:
                val = float(sz_match.group(1))
                unit = sz_match.group(2)
                if unit == "em":
                    info["font_size"] = val * 16
                elif unit in ("px", "pt", ";"):
                    info["font_size"] = val

            if "text-align:center" in css_text or "text-align: center" in css_text:
                info["text_align"] = "center"

            props[cls_name] = info

        return props

    def _get_class_info(self, element, class_props: dict) -> dict:
        cls = element.get("class", "")
        if cls and cls in class_props:
            return class_props[cls]
        return {"bold": False, "italic": False, "font_size": None, "text_align": None}

    def _process_element(self, element, tokens, class_props, image_dir, extract_images):
        for child in element:
            tag = child.tag

            if tag == "style":
                continue
            elif tag == "table":
                self._process_table(child, tokens)
            elif tag == "img":
                if extract_images:
                    src = child.get("src", "")
                    alt = child.get("alt", "")
                    if src and not src.startswith("x-apple") and not src.startswith("data:"):
                        tokens.append(DocToken(kind="image", image_path=src, caption=alt))
            elif tag == "p":
                self._process_paragraph(child, tokens, class_props)
            elif tag in ("div", "section", "article", "main", "header", "footer"):
                self._process_element(child, tokens, class_props, image_dir, extract_images)
            elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
                level = int(tag[1])
                text = child.text_content().strip()
                if text:
                    tokens.append(DocToken(kind="heading", text=text, level=level))
            elif tag in ("ul", "ol"):
                self._process_list(child, tokens, class_props, ordered=(tag == "ol"), level=0)

    def _process_paragraph(self, p_element, tokens, class_props):
        text_parts = []
        max_font_size = 0
        has_bold = False

        for span in p_element.iter("span"):
            txt = span.text or ""
            tail = span.tail or ""
            cls_info = self._get_class_info(span, class_props)

            if cls_info["font_size"] and cls_info["font_size"] > max_font_size:
                max_font_size = cls_info["font_size"]
            if cls_info["bold"]:
                has_bold = True

            if txt.strip():
                if cls_info["bold"] and cls_info["italic"]:
                    text_parts.append(f"***{txt.strip()}***")
                elif cls_info["bold"]:
                    text_parts.append(f"**{txt.strip()}**")
                elif cls_info["italic"]:
                    text_parts.append(f"*{txt.strip()}*")
                else:
                    text_parts.append(txt)
            elif txt:
                text_parts.append(txt)

            if tail.strip():
                text_parts.append(tail)
            elif tail:
                text_parts.append(tail)

        if not text_parts:
            direct = (p_element.text or "").strip()
            if direct:
                text_parts.append(direct)

        full_text = "".join(text_parts).strip()
        full_text = re.sub(r'[\xa0\u200b]+', ' ', full_text)
        full_text = re.sub(r'\s+', ' ', full_text).strip()

        if not full_text:
            return

        h1_thresh = getattr(self, '_h1_thresh', 28)
        h2_thresh = getattr(self, '_h2_thresh', 22)

        if max_font_size >= h1_thresh:
            clean = re.sub(r'\*+', '', full_text).strip()
            tokens.append(DocToken(kind="heading", text=clean, level=1))
        elif max_font_size >= h2_thresh:
            clean = re.sub(r'\*+', '', full_text).strip()
            tokens.append(DocToken(kind="heading", text=clean, level=2))
        elif has_bold and len(full_text) < 120:
            clean = re.sub(r'\*+', '', full_text).strip()
            tokens.append(DocToken(kind="heading", text=clean, level=3))
        else:
            tokens.append(DocToken(kind="paragraph", text=full_text))

    def _process_table(self, table_element, tokens):
        rows_data = []
        for tr in table_element.findall(".//tr"):
            cells = []
            for td in tr.findall(".//td"):
                text = td.text_content().strip()
                text = re.sub(r'[\xa0\u200b]+', ' ', text)
                text = re.sub(r'\s+', ' ', text).strip()
                cells.append(text)
            if cells and any(c for c in cells):
                rows_data.append(cells)

        if not rows_data:
            return

        headers = rows_data[0]
        body = rows_data[1:] if len(rows_data) > 1 else []
        tokens.append(DocToken(kind="table", headers=headers, rows=body))

    def _process_list(self, list_element, tokens, class_props, ordered=False, level=0):
        for li in list_element.findall("li"):
            text_parts = []
            for span in li.iter("span"):
                t = span.text or ""
                if t.strip():
                    cls_info = self._get_class_info(span, class_props)
                    if cls_info["bold"]:
                        text_parts.append(f"**{t.strip()}**")
                    elif cls_info["italic"]:
                        text_parts.append(f"*{t.strip()}*")
                    else:
                        text_parts.append(t)
            if text_parts:
                text = "".join(text_parts).strip()
            else:
                text = li.text_content().strip()

            text = re.sub(r'[\xa0\u200b]+', ' ', text)
            text = re.sub(r'\s+', ' ', text).strip()

            if text:
                tokens.append(DocToken(kind="list", items=[(level, text, ordered)]))

            for sub in li.findall("ul"):
                self._process_list(sub, tokens, class_props, ordered=False, level=level + 1)
            for sub in li.findall("ol"):
                self._process_list(sub, tokens, class_props, ordered=True, level=level + 1)


# ---------------------------------------------------------------------------
# DocxConverter (#3, #8)
# ---------------------------------------------------------------------------

class DocxConverter:
    """Converts DOCX → list of DocToken."""

    def convert(self, docx_path: str, image_dir: str, extract_images: bool = True,
                verbose: bool = False) -> Tuple[List[DocToken], dict]:
        # Detect HTML disguised as .docx
        try:
            with open(docx_path, "rb") as f:
                header = f.read(20)
            if header.lstrip().startswith(b"<"):
                if verbose:
                    print(f"  (HTML file disguised as .docx)")
                return HtmlDocConverter().convert(docx_path, image_dir, extract_images, verbose)
        except Exception:
            pass

        try:
            from docx import Document
            from docx.oxml.ns import qn
        except ImportError:
            print("ERROR: python-docx is not installed. pip install python-docx", file=sys.stderr)
            sys.exit(1)

        document = Document(docx_path)
        tokens: List[DocToken] = []
        meta = {}

        # --- Extended metadata (#8) ---
        props = document.core_properties
        if props.title:
            meta["title"] = props.title
        if props.author:
            meta["author"] = props.author
        if props.subject:
            meta["subject"] = props.subject
        if props.keywords:
            meta["keywords"] = props.keywords
        if props.category:
            meta["category"] = props.category
        if props.created:
            try:
                meta["date"] = props.created.strftime("%Y-%m-%d")
            except Exception:
                pass

        # Extract images from relationships
        img_counter = 0
        image_map = {}  # rId → filename
        if extract_images:
            for rel in document.part.rels.values():
                if "image" in rel.reltype:
                    if rel.is_external:
                        continue
                    try:
                        img_counter += 1
                        ct = getattr(rel.target_part, "content_type", "image/png")
                        ext = ct.split("/")[-1].replace("jpeg", "jpg")
                        if ext not in ("png", "jpg", "gif", "bmp", "svg", "tiff", "webp"):
                            ext = "png"
                        img_name = f"docx_img_{img_counter}.{ext}"
                        img_path = os.path.join(image_dir, img_name)
                        with open(img_path, "wb") as f:
                            f.write(rel.target_part.blob)
                        image_map[rel.rId] = img_name
                    except (ValueError, AttributeError):
                        img_counter -= 1

        used_images = set()

        # Process document body elements in order
        for element in document.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                self._process_paragraph(element, document, tokens, image_map, used_images, qn)
            elif tag == "tbl":
                self._process_table_element(element, document, tokens, qn)

        # Add any remaining unreferenced images
        for rId, img_name in image_map.items():
            if rId not in used_images:
                tokens.append(DocToken(kind="image", image_path=img_name))

        tokens = _merge_lists(tokens)
        return tokens, meta

    def _process_paragraph(self, element, document, tokens, image_map, used_images, qn):
        from docx.text.paragraph import Paragraph
        para = Paragraph(element, document)
        style_name = para.style.name if para.style else ""

        # Check for inline images in runs
        for run in para.runs:
            drawings = run.element.findall(f".//{qn('wp:inline')}")
            drawings += run.element.findall(f".//{qn('wp:anchor')}")
            for drawing in drawings:
                blip = drawing.find(f".//{qn('a:blip')}")
                if blip is not None:
                    embed = blip.get(qn("r:embed"))
                    if embed and embed in image_map:
                        used_images.add(embed)
                        tokens.append(DocToken(kind="image", image_path=image_map[embed]))

        # Heading
        if style_name.startswith("Heading") or style_name == "Title" or style_name == "Subtitle":
            if style_name == "Title":
                level = 1
            elif style_name == "Subtitle":
                if para.text.strip():
                    tokens.append(DocToken(kind="paragraph", text=f"*{para.text.strip()}*"))
                return
            else:
                try:
                    level = int(style_name.split()[-1])
                except (ValueError, IndexError):
                    level = 1
            level = max(1, min(level, 4))
            tokens.append(DocToken(kind="heading", text=para.text.strip(), level=level))
            return

        # Quote
        if "Quote" in style_name or "Block Text" in style_name or "Intense Quote" in style_name:
            if para.text.strip():
                tokens.append(DocToken(kind="quote", text=para.text.strip()))
            return

        # Code style (#7 — code detection by style)
        if any(kw in style_name.lower() for kw in ("code", "source", "listing", "preformatted",
                                                      "html preformatted")):
            if para.text.strip():
                tokens.append(DocToken(kind="code", text=para.text))
            return

        # List (by XML numPr or by style name) (#3 — improved nested list handling)
        is_list = False
        indent_level = 0
        ordered = False

        numPr = element.find(f".//{qn('w:numPr')}")
        if numPr is not None:
            is_list = True
            ilvl_el = numPr.find(qn("w:ilvl"))
            numId_el = numPr.find(qn("w:numId"))
            indent_level = int(ilvl_el.get(qn("w:val"), "0")) if ilvl_el is not None else 0
            num_id = int(numId_el.get(qn("w:val"), "0")) if numId_el is not None else 0
            ordered = self._is_ordered_list(document, num_id, indent_level)
        elif "List" in style_name:
            is_list = True
            ordered = "Number" in style_name
            try:
                indent = para.paragraph_format.left_indent
                if indent:
                    from docx.shared import Inches
                    indent_level = max(0, int(indent / Inches(0.25)) - 1)
            except Exception:
                pass

        if is_list:
            text = self._extract_formatted_text(para)
            if text.strip():
                tokens.append(DocToken(kind="list", items=[(indent_level, text.strip(), ordered)]))
            return

        # Regular paragraph — detect monospace font (code)
        text = self._extract_formatted_text(para)
        if text.strip():
            # Check if all runs use monospace font → code block
            if para.runs and all(self._is_mono_font(r) for r in para.runs if r.text.strip()):
                tokens.append(DocToken(kind="code", text=para.text))
            else:
                tokens.append(DocToken(kind="paragraph", text=text.strip()))

    @staticmethod
    def _is_mono_font(run) -> bool:
        """Checks whether the run uses a monospace font."""
        font = run.font
        if font.name:
            name = font.name.lower()
            return any(kw in name for kw in (
                "courier", "mono", "consolas", "menlo", "source code",
                "fira code", "inconsolata", "dejavu sans mono",
                "liberation mono", "roboto mono", "jetbrains", "hack",
            ))
        return False

    def _extract_formatted_text(self, para) -> str:
        parts = []
        for run in para.runs:
            t = run.text
            if not t:
                continue
            font = run.font
            is_bold = run.bold or False
            is_italic = run.italic or False
            is_code = False
            if font.name and any(m in font.name.lower()
                                  for m in ("courier", "mono", "consolas", "menlo")):
                is_code = True
            is_super = font.superscript or False
            is_sub = font.subscript or False
            is_strike = font.strike or False
            is_underline = font.underline or False

            if is_code:
                parts.append(f"`{t}`")
            elif is_bold and is_italic:
                parts.append(f"***{t}***")
            elif is_bold:
                parts.append(f"**{t}**")
            elif is_italic:
                parts.append(f"*{t}*")
            elif is_strike:
                parts.append(f"~~{t}~~")
            elif is_underline and not is_bold:
                # Underline without bold — underline (no MD equivalent, use italic)
                parts.append(f"*{t}*")
            elif is_super:
                parts.append(f"^{t}^")
            elif is_sub:
                parts.append(f"~{t}~")
            else:
                parts.append(t)
        return "".join(parts)

    def _is_ordered_list(self, document, num_id: int, ilvl: int = 0) -> bool:
        """Determines list type from the numbering definition with nesting level (#3)."""
        try:
            numbering = document.part.numbering_part.numbering_definitions
            from docx.oxml.ns import qn

            # Find abstractNumId for this numId
            abstract_num_id = None
            for num_el in numbering.findall(qn("w:num")):
                if num_el.get(qn("w:numId")) == str(num_id):
                    abstract_ref = num_el.find(qn("w:abstractNumId"))
                    if abstract_ref is not None:
                        abstract_num_id = abstract_ref.get(qn("w:val"))
                    break

            # Find format for the specific level
            for abstract in numbering.findall(qn("w:abstractNum")):
                aid = abstract.get(qn("w:abstractNumId"))
                if abstract_num_id is not None and aid != abstract_num_id:
                    continue
                for lvl in abstract.findall(qn("w:lvl")):
                    lvl_val = lvl.get(qn("w:ilvl"), "0")
                    if int(lvl_val) == ilvl:
                        num_fmt = lvl.find(qn("w:numFmt"))
                        if num_fmt is not None:
                            fmt = num_fmt.get(qn("w:val"), "")
                            if fmt in ("decimal", "lowerLetter", "upperLetter",
                                       "lowerRoman", "upperRoman"):
                                return True
                            elif fmt == "bullet":
                                return False
                if abstract_num_id is not None:
                    break
        except Exception:
            pass
        return False

    def _process_table_element(self, element, document, tokens, qn):
        from docx.table import Table
        table = Table(element, document)
        if not table.rows:
            return
        headers = [cell.text.strip() for cell in table.rows[0].cells]
        rows_data = []
        for row in table.rows[1:]:
            rows_data.append([cell.text.strip() for cell in row.cells])
        tokens.append(DocToken(kind="table", headers=headers, rows=rows_data))


# ---------------------------------------------------------------------------
# PptxConverter (#4, #5)
# ---------------------------------------------------------------------------

class PptxConverter:
    """Converts PPTX → list of DocToken."""

    CHART_TYPE_MAP = {}  # populated in convert()

    def convert(self, pptx_path: str, image_dir: str, extract_images: bool = True,
                verbose: bool = False) -> Tuple[List[DocToken], dict]:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.enum.chart import XL_CHART_TYPE
        except ImportError:
            print("ERROR: python-pptx is not installed. pip install python-pptx", file=sys.stderr)
            sys.exit(1)

        self.CHART_TYPE_MAP = {
            XL_CHART_TYPE.PIE: "pie",
            XL_CHART_TYPE.PIE_EXPLODED: "pie",
            XL_CHART_TYPE.DOUGHNUT: "pie",
            XL_CHART_TYPE.DOUGHNUT_EXPLODED: "pie",
            XL_CHART_TYPE.BAR_CLUSTERED: "bar",
            XL_CHART_TYPE.BAR_STACKED: "bar",
            XL_CHART_TYPE.BAR_STACKED_100: "bar",
            XL_CHART_TYPE.COLUMN_CLUSTERED: "bar",
            XL_CHART_TYPE.COLUMN_STACKED: "bar",
            XL_CHART_TYPE.COLUMN_STACKED_100: "bar",
            XL_CHART_TYPE.LINE: "line",
            XL_CHART_TYPE.LINE_MARKERS: "line",
            XL_CHART_TYPE.LINE_STACKED: "line",
            XL_CHART_TYPE.LINE_MARKERS_STACKED: "line",
            XL_CHART_TYPE.AREA: "area",
            XL_CHART_TYPE.AREA_STACKED: "area",
            XL_CHART_TYPE.AREA_STACKED_100: "area",
        }

        prs = Presentation(pptx_path)
        tokens: List[DocToken] = []
        meta = {}

        # --- Extended metadata (#8) ---
        props = prs.core_properties
        if props.title:
            meta["title"] = props.title
        if props.author:
            meta["author"] = props.author
        if props.subject:
            meta["subject"] = props.subject
        if props.keywords:
            meta["keywords"] = props.keywords
        if props.category:
            meta["category"] = props.category
        if props.created:
            try:
                meta["date"] = props.created.strftime("%Y-%m-%d")
            except Exception:
                pass

        img_counter = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            title_text = self._get_slide_title(slide)
            if title_text:
                tokens.append(DocToken(kind="heading", text=title_text, level=2))
            else:
                tokens.append(DocToken(kind="heading", text=f"Slide {slide_num}", level=2))

            shapes = sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0))
            for shape in shapes:
                img_counter = self._process_shape(
                    shape, tokens, image_dir, extract_images, img_counter,
                    slide_num, verbose, title_text
                )

            # Slide notes → footnote
            if slide.has_notes_slide:
                try:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        tokens.append(DocToken(kind="footnote",
                                               text=f"Slide {slide_num}: {notes_text}"))
                except Exception:
                    pass

        tokens = _merge_lists(tokens)
        return tokens, meta

    @staticmethod
    def _get_slide_title(slide) -> str:
        if slide.shapes.title:
            try:
                return slide.shapes.title.text.strip()
            except AttributeError:
                pass
        for shape in slide.placeholders:
            if shape.placeholder_format.idx in (0, 1):
                try:
                    txt = shape.text.strip()
                    if txt:
                        return txt
                except AttributeError:
                    pass
        return ""

    def _process_shape(self, shape, tokens, image_dir, extract_images,
                       img_counter, slide_num, verbose, title_text):
        from pptx.util import Emu

        # Skip title placeholder
        try:
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                if shape.placeholder_format.idx in (0, 1) and hasattr(shape, "text"):
                    if shape.text.strip() == title_text:
                        return img_counter
        except Exception:
            pass

        # Group shape → recurse
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                for child in shape.shapes:
                    img_counter = self._process_shape(
                        child, tokens, image_dir, extract_images,
                        img_counter, slide_num, verbose, title_text
                    )
            except Exception:
                pass
            return img_counter

        # Table
        if shape.has_table:
            tbl = shape.table
            all_rows = list(tbl.rows)
            if all_rows:
                headers = [cell.text.strip() for cell in all_rows[0].cells]
                rows_data = []
                for row in all_rows[1:]:
                    rows_data.append([cell.text.strip() for cell in row.cells])
                tokens.append(DocToken(kind="table", headers=headers, rows=rows_data))
            return img_counter

        # Chart
        if shape.has_chart:
            chart_token = self._extract_chart(shape.chart, image_dir, img_counter, slide_num)
            if chart_token:
                tokens.append(chart_token)
                if chart_token.kind == "image":
                    img_counter += 1
            return img_counter

        # Image
        if hasattr(shape, "image") and extract_images:
            try:
                blob = shape.image.blob
                ct = shape.image.content_type
                ext = ct.split("/")[-1].replace("jpeg", "jpg") if ct else "png"
                if ext not in ("png", "jpg", "gif", "bmp", "svg", "tiff", "webp"):
                    ext = "png"
                img_counter += 1
                img_name = f"slide_{slide_num}_img_{img_counter}.{ext}"
                img_path = os.path.join(image_dir, img_name)
                with open(img_path, "wb") as f:
                    f.write(blob)
                alt = shape.name or ""
                tokens.append(DocToken(kind="image", image_path=img_name, caption=alt))
                return img_counter
            except Exception:
                pass

        # Text frame (#4 — extended formatting)
        if shape.has_text_frame:
            paras = list(shape.text_frame.paragraphs)
            is_body_ph = False
            try:
                if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                    is_body_ph = shape.placeholder_format.idx == 1
            except Exception:
                pass

            has_levels = any((p.level or 0) > 0 for p in paras)
            multi_para = sum(1 for p in paras if p.text.strip()) > 1

            for para in paras:
                text = self._extract_pptx_formatted_text(para)
                if not text.strip():
                    continue

                level = para.level or 0
                is_bullet = self._is_bullet(para)

                if level > 0 or is_bullet or (is_body_ph and multi_para) or \
                   (has_levels and multi_para):
                    ordered = self._is_numbered(para)
                    tokens.append(DocToken(kind="list",
                                           items=[(level, text.strip(), ordered)]))
                else:
                    tokens.append(DocToken(kind="paragraph", text=text.strip()))

        return img_counter

    def _extract_pptx_formatted_text(self, para) -> str:
        """Extracts formatted text from a PPTX paragraph (#4)."""
        parts = []
        for run in para.runs:
            t = run.text
            if not t:
                continue
            font = run.font

            is_bold = font.bold or False
            is_italic = font.italic or False
            is_underline = font.underline or False
            is_strike = getattr(font, 'strikethrough', None) or False

            # Code detection by font (#4)
            is_code = False
            font_name = getattr(font, 'name', None) or ""
            if font_name:
                name_lower = font_name.lower()
                if any(kw in name_lower for kw in (
                    "courier", "mono", "consolas", "menlo", "source code",
                    "fira code", "inconsolata",
                )):
                    is_code = True

            if is_code:
                parts.append(f"`{t}`")
            elif is_bold and is_italic:
                parts.append(f"***{t}***")
            elif is_bold:
                parts.append(f"**{t}**")
            elif is_italic:
                parts.append(f"*{t}*")
            elif is_strike:
                parts.append(f"~~{t}~~")
            elif is_underline and not is_bold:
                parts.append(f"*{t}*")
            else:
                parts.append(t)

        # Hyperlinks (#4)
        try:
            for run in para.runs:
                if run.hyperlink and run.hyperlink.address:
                    # Replace run text with a markdown link
                    link_text = run.text or run.hyperlink.address
                    idx = None
                    for i, p in enumerate(parts):
                        if link_text in p:
                            idx = i
                            break
                    if idx is not None:
                        parts[idx] = f"[{link_text}]({run.hyperlink.address})"
        except Exception:
            pass

        return "".join(parts)

    def _is_bullet(self, para) -> bool:
        pPr = para._p.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        if pPr is not None:
            buNone = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buNone")
            if buNone is not None:
                return False
            buChar = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar")
            buAutoNum = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
            if buChar is not None or buAutoNum is not None:
                return True
        return False

    def _is_numbered(self, para) -> bool:
        pPr = para._p.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
        if pPr is not None:
            buAutoNum = pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum")
            return buAutoNum is not None
        return False

    def _extract_chart(self, chart, image_dir, img_counter, slide_num) -> Optional[DocToken]:
        """Extracts a chart → ```chart or image fallback (#5)."""
        try:
            chart_type_enum = chart.chart_type
            md_type = self.CHART_TYPE_MAP.get(chart_type_enum)

            title = ""
            if chart.has_title:
                title = chart.chart_title.text_frame.text.strip()

            plot = chart.plots[0]

            categories = []
            try:
                categories = [str(c) for c in plot.categories]
            except Exception:
                pass

            series_list = []
            for series in plot.series:
                name = series.name if hasattr(series, 'name') and series.name else "Series"
                values = []
                try:
                    values = [v for v in series.values]
                except Exception:
                    pass
                if values:
                    series_list.append((name, values))

            if not series_list:
                # No data → try rendering as image (#5)
                return self._render_chart_as_image(
                    chart, title, categories, series_list,
                    image_dir, img_counter, slide_num
                )

            if md_type is None:
                # Unknown chart type → render via matplotlib (#5)
                img_token = self._render_chart_as_image(
                    chart, title, categories, series_list,
                    image_dir, img_counter, slide_num
                )
                if img_token:
                    return img_token
                # Fallback: use bar as the default type
                md_type = "bar"

            # Build ```chart spec
            lines = [f"type: {md_type}"]
            if title:
                lines.append(f"title: {title}")
            if categories and md_type != "pie":
                lines.append(f"labels: {', '.join(categories)}")

            if md_type == "pie":
                ser_name, ser_vals = series_list[0]
                if categories:
                    for cat, val in zip(categories, ser_vals):
                        v = int(val) if val == int(val) else val
                        lines.append(f"{cat}: {v}")
                else:
                    for i, val in enumerate(ser_vals):
                        v = int(val) if val == int(val) else val
                        lines.append(f"Item {i + 1}: {v}")
            else:
                for name, vals in series_list:
                    formatted = ", ".join(
                        str(int(v)) if v == int(v) else str(v)
                        for v in vals
                    )
                    lines.append(f"{name}: {formatted}")

            chart_spec = "\n".join(lines)
            return DocToken(kind="chart", chart_spec=chart_spec)

        except Exception:
            return None

    @staticmethod
    def _render_chart_as_image(chart, title: str, categories: list,
                               series_list: list, image_dir: str,
                               img_counter: int, slide_num: int) -> Optional[DocToken]:
        """Renders a chart as an image via matplotlib (#5)."""
        try:
            import matplotlib
            matplotlib.use('Agg')
            import matplotlib.pyplot as plt
        except ImportError:
            return None

        if not series_list:
            return None

        try:
            fig, ax = plt.subplots(figsize=(8, 5))

            labels = categories if categories else [str(i + 1)
                                                     for i in range(len(series_list[0][1]))]

            if len(series_list) == 1:
                name, vals = series_list[0]
                ax.bar(labels[:len(vals)], vals, label=name)
            else:
                import numpy as np
                n = len(series_list)
                width = 0.8 / n
                x = np.arange(len(labels))
                for idx, (name, vals) in enumerate(series_list):
                    offset = (idx - n / 2 + 0.5) * width
                    ax.bar(x[:len(vals)] + offset, vals, width=width, label=name)
                ax.set_xticks(x)
                ax.set_xticklabels(labels, rotation=45, ha='right')

            if title:
                ax.set_title(title)
            if len(series_list) > 1:
                ax.legend()

            plt.tight_layout()
            img_counter += 1
            img_name = f"chart_{slide_num}_{img_counter}.png"
            img_path = os.path.join(image_dir, img_name)
            fig.savefig(img_path, dpi=150, bbox_inches='tight')
            plt.close(fig)

            return DocToken(kind="image", image_path=img_name,
                            caption=title or "Chart")
        except Exception:
            return None


# ---------------------------------------------------------------------------
# MarkdownRenderer (#7, #8)
# ---------------------------------------------------------------------------

class MarkdownRenderer:
    """Renders a list of DocToken into Markdown text."""

    def render(self, tokens: List[DocToken], meta: dict, image_dir_name: str = "images") -> str:
        parts = []

        # YAML front matter (#8 — extended metadata)
        fm_lines = []
        if meta.get("title"):
            fm_lines.append(f"title: \"{meta['title']}\"")
        if meta.get("author"):
            fm_lines.append(f"author: \"{meta['author']}\"")
        if meta.get("date"):
            fm_lines.append(f"date: \"{meta['date']}\"")
        if meta.get("subject"):
            fm_lines.append(f"subject: \"{meta['subject']}\"")
        if meta.get("keywords"):
            fm_lines.append(f"keywords: \"{meta['keywords']}\"")
        if meta.get("category"):
            fm_lines.append(f"category: \"{meta['category']}\"")
        if meta.get("creator"):
            fm_lines.append(f"creator: \"{meta['creator']}\"")
        if fm_lines:
            parts.append("---")
            parts.extend(fm_lines)
            parts.append("---")
            parts.append("")

        title_emitted = False
        footnotes = []

        for tok in tokens:
            if tok.kind == "heading":
                prefix = "#" * tok.level
                parts.append(f"{prefix} {tok.text}")
                parts.append("")
                if tok.level == 1 and not title_emitted:
                    title_emitted = True

            elif tok.kind == "paragraph":
                parts.append(tok.text)
                parts.append("")

            elif tok.kind == "quote":
                for line in tok.text.split("\n"):
                    parts.append(f"> {line}")
                parts.append("")

            elif tok.kind == "table":
                self._render_table(parts, tok)

            elif tok.kind == "image":
                caption = tok.caption or ""
                path = f"{image_dir_name}/{tok.image_path}" if tok.image_path else ""
                parts.append(f"![{caption}]({path})")
                parts.append("")

            elif tok.kind == "code":
                # Auto-detect language (#7)
                lang = tok.lang or detect_code_language(tok.text)
                parts.append(f"```{lang}")
                parts.append(tok.text)
                parts.append("```")
                parts.append("")

            elif tok.kind == "list":
                self._render_list(parts, tok)
                parts.append("")

            elif tok.kind == "chart":
                parts.append("```chart")
                parts.append(tok.chart_spec)
                parts.append("```")
                parts.append("")

            elif tok.kind == "hr":
                parts.append("---")
                parts.append("")

            elif tok.kind == "footnote":
                footnotes.append(tok.text)

        # Footnotes at end
        if footnotes:
            parts.append("")
            for i, fn in enumerate(footnotes, 1):
                parts.append(f"[^{i}]: {fn}")
            parts.append("")

        text = "\n".join(parts)
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip() + "\n"

    def _render_table(self, parts: list, tok: DocToken):
        headers = tok.headers
        rows = tok.rows

        if not headers and not rows:
            return

        if tok.caption:
            parts.append(f"Table: {tok.caption}")
            parts.append("")

        if not headers and rows:
            headers = rows[0]
            rows = rows[1:]

        col_count = len(headers)
        for row in rows:
            col_count = max(col_count, len(row))

        while len(headers) < col_count:
            headers.append("")

        parts.append("| " + " | ".join(
            h.replace("|", "\\|").replace("\n", " ") for h in headers
        ) + " |")

        aligns = tok.aligns
        sep_parts = []
        for i in range(col_count):
            if aligns and i < len(aligns):
                a = aligns[i]
                if a == "center":
                    sep_parts.append(":---:")
                elif a == "right":
                    sep_parts.append("---:")
                else:
                    sep_parts.append("---")
            else:
                sep_parts.append("---")
        parts.append("| " + " | ".join(sep_parts) + " |")

        for row in rows:
            padded = list(row) + [""] * (col_count - len(row))
            parts.append("| " + " | ".join(
                c.replace("|", "\\|").replace("\n", " ") for c in padded
            ) + " |")

        parts.append("")

    def _render_list(self, parts: list, tok: DocToken):
        ordered_counters = {}
        for level, text, ordered in tok.items:
            indent = "  " * level
            if ordered:
                counter = ordered_counters.get(level, 0) + 1
                ordered_counters[level] = counter
                parts.append(f"{indent}{counter}. {text}")
            else:
                parts.append(f"{indent}- {text}")
                ordered_counters.pop(level, None)


# ---------------------------------------------------------------------------
# MarkdownLinter — post-processing (#6)
# ---------------------------------------------------------------------------

class MarkdownLinter:
    """Normalizes and cleans Markdown output."""

    @staticmethod
    def lint(text: str) -> str:
        lines = text.split('\n')

        # 1. Strip trailing whitespace
        lines = [line.rstrip() for line in lines]
        text = '\n'.join(lines)

        # 2. Merge adjacent bold: **text** **more** → **text more**
        text = re.sub(r'\*\*([^*]+)\*\*\s+\*\*([^*]+)\*\*', r'**\1 \2**', text)

        # 3. Merge adjacent italic: *text* *more* → *text more*
        text = re.sub(r'(?<!\*)\*([^*]+)\*\s+\*([^*]+)\*(?!\*)', r'*\1 \2*', text)

        # 4. Drop empty markers: ****  or ** **
        text = re.sub(r'(?<!\w)\*{2,3}\s*\*{2,3}(?!\w)', '', text)
        text = re.sub(r'(?<!\w)\*\s+\*(?!\w)', ' ', text)

        # 5. Blank line before headings (if missing)
        text = re.sub(r'([^\n])\n(#{1,6} )', r'\1\n\n\2', text)

        # 6. Blank line after headings (if missing)
        text = re.sub(r'(#{1,6} [^\n]+)\n([^\n#])', r'\1\n\n\2', text)

        # 7. No more than 2 consecutive blank lines
        text = re.sub(r'\n{3,}', '\n\n', text)

        # 8. Trim spaces inside inline markers: ** text ** → **text**
        text = re.sub(r'\*\*\s+([^*]+?)\s+\*\*', r'**\1**', text)
        text = re.sub(r'(?<!\*)\*\s+([^*]+?)\s+\*(?!\*)', r'*\1*', text)

        # 9. Drop duplicated list markers: - - text → - text
        text = re.sub(r'^(\s*)- - ', r'\1- ', text, flags=re.M)

        # 10. Normalize tabs in list indentation → 2 spaces
        def fix_list_indent(m):
            tabs = m.group(1).count('\t')
            spaces = ' ' * (tabs * 2)
            rest = m.group(1).replace('\t', '  ')
            return rest + m.group(2)
        text = re.sub(r'^(\t+)([-*\d])', fix_list_indent, text, flags=re.M)

        # 11. File must end with a single newline
        text = text.rstrip('\n') + '\n'

        return text


# ---------------------------------------------------------------------------
# Legacy format conversion (.ppt → .pptx, .doc → .docx)
# ---------------------------------------------------------------------------

def _convert_legacy_format(input_path: str, ext: str, verbose: bool = False) -> Optional[str]:
    """Convert old binary .ppt/.doc to .pptx/.docx via system tools."""
    import subprocess
    import tempfile

    basename = os.path.basename(input_path)

    if ext == ".doc":
        tmp_dir = tempfile.mkdtemp(prefix="doc2md_")
        out_file = os.path.join(tmp_dir, os.path.splitext(basename)[0] + ".docx")
        try:
            subprocess.run(
                ["textutil", "-convert", "docx", "-output", out_file, input_path],
                check=True, capture_output=True, text=True,
            )
            if verbose:
                print(f"  .doc → .docx: {out_file}")
            return out_file
        except FileNotFoundError:
            print(f"ERROR: {basename}: textutil not found (macOS required)", file=sys.stderr)
            return None
        except subprocess.CalledProcessError as e:
            print(f"ERROR: {basename}: textutil: {e.stderr.strip()}", file=sys.stderr)
            return None

    if ext == ".ppt":
        soffice = None
        for path in [
            "/usr/bin/soffice",
            "/usr/local/bin/soffice",
            "/opt/homebrew/bin/soffice",
            "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        ]:
            if os.path.isfile(path):
                soffice = path
                break
        if soffice is None:
            import shutil
            soffice = shutil.which("soffice")

        if soffice is None:
            print(f"ERROR: {basename}: .ppt format requires LibreOffice for conversion. "
                  f"Install: brew install --cask libreoffice", file=sys.stderr)
            return None

        tmp_dir = tempfile.mkdtemp(prefix="doc2md_")
        try:
            subprocess.run(
                [soffice, "--headless", "--convert-to", "pptx",
                 "--outdir", tmp_dir, input_path],
                check=True, capture_output=True, text=True,
            )
            out_file = os.path.join(tmp_dir, os.path.splitext(basename)[0] + ".pptx")
            if os.path.isfile(out_file):
                if verbose:
                    print(f"  .ppt → .pptx: {out_file}")
                return out_file
            print(f"ERROR: {basename}: soffice did not create .pptx", file=sys.stderr)
            return None
        except subprocess.CalledProcessError as e:
            print(f"ERROR: {basename}: soffice: {e.stderr.strip()}", file=sys.stderr)
            return None

    return None


# ---------------------------------------------------------------------------
# Single-file conversion (for sequential and parallel modes) (#11)
# ---------------------------------------------------------------------------

def _convert_one_file(work_item: tuple) -> tuple:
    """Converts one file. Returns (success: bool, message: str, tmp_file: str|None)."""
    (input_path, output_path, image_dir_name, extract_images,
     verbose, force, converters_map) = work_item

    input_path = os.path.abspath(input_path)
    original_path = input_path
    tmp_file = None

    if not os.path.isfile(input_path):
        return (False, f"ERROR: file not found: {input_path}", None)

    ext = os.path.splitext(input_path)[1].lower()
    converter_key = ext
    if converter_key not in converters_map:
        return (False, f"ERROR: unsupported format: {ext} ({input_path})", None)

    # Legacy formats
    if ext in (".ppt", ".doc"):
        converted = _convert_legacy_format(input_path, ext, verbose=verbose)
        if converted is None:
            return (False, f"ERROR: could not convert {ext}: {input_path}", None)
        tmp_file = converted
        input_path = converted
        ext = os.path.splitext(converted)[1].lower()
        converter_key = ext

    # Output path
    if output_path:
        out_path = os.path.abspath(output_path)
    else:
        out_path = os.path.splitext(original_path)[0] + ".md"

    if os.path.exists(out_path) and not force:
        return (False, f"SKIP: {out_path} already exists (use --force)", tmp_file)

    out_dir = os.path.dirname(out_path)
    image_dir = os.path.join(out_dir, image_dir_name)
    if extract_images:
        os.makedirs(image_dir, exist_ok=True)

    if verbose:
        print(f"Converting: {input_path} → {out_path}")

    try:
        converter_cls_name = converters_map[converter_key]
        # Build a converter instance by class name
        converter_cls = {
            "PDFConverter": PDFConverter,
            "DocxConverter": DocxConverter,
            "PptxConverter": PptxConverter,
        }[converter_cls_name]

        converter = converter_cls()
        tokens, meta = converter.convert(
            input_path, image_dir,
            extract_images=extract_images,
            verbose=verbose,
        )

        renderer = MarkdownRenderer()
        md_text = renderer.render(tokens, meta, image_dir_name=image_dir_name)

        # Markdown lint (#6)
        md_text = MarkdownLinter.lint(md_text)

        with open(out_path, "w", encoding="utf-8") as f:
            f.write(md_text)

        basename = os.path.basename(out_path)
        tok_count = len(tokens)
        img_count = sum(1 for t in tokens if t.kind == "image")
        tbl_count = sum(1 for t in tokens if t.kind == "table")
        chart_count = sum(1 for t in tokens if t.kind == "chart")
        code_count = sum(1 for t in tokens if t.kind == "code")

        info_parts = [f"{basename}"]
        if tok_count:
            info_parts.append(f"{tok_count} elements")
        if img_count:
            info_parts.append(f"{img_count} images")
        if tbl_count:
            info_parts.append(f"{tbl_count} tables")
        if chart_count:
            info_parts.append(f"{chart_count} charts")
        if code_count:
            info_parts.append(f"{code_count} code")

        return (True, f"  OK: {', '.join(info_parts)}", tmp_file)

    except Exception as e:
        import traceback
        msg = f"ERROR: {os.path.basename(input_path)}: {e}"
        if verbose:
            msg += "\n" + traceback.format_exc()
        return (False, msg, tmp_file)


# ---------------------------------------------------------------------------
# CLI & main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="doc2md — PDF / DOCX / PPTX → Markdown converter",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""Examples:
  python3 doc2md.py lecture.pdf
  python3 doc2md.py report.docx -o report.md
  python3 doc2md.py slides.pptx --image-dir assets
  python3 doc2md.py *.pptx --verbose
  python3 doc2md.py doc.pdf doc2.docx slides.pptx  # batch
  python3 doc2md.py *.pdf --parallel  # parallel processing
""",
    )
    parser.add_argument("input", nargs="+", help="Input files (PDF, DOCX, PPTX)")
    parser.add_argument("-o", "--output",
                        help="Output .md file (only for a single input file)")
    parser.add_argument("--image-dir", default="images",
                        help="Directory for extracted images (default: images)")
    parser.add_argument("--no-images", action="store_true",
                        help="Do not extract images")
    parser.add_argument("--verbose", "-v", action="store_true",
                        help="Verbose output")
    parser.add_argument("--force", action="store_true",
                        help="Overwrite existing .md file")
    parser.add_argument("--parallel", "-j", action="store_true",
                        help="Parallel processing of multiple files (#11)")

    args = parser.parse_args()

    if args.output and len(args.input) > 1:
        parser.error("--output can only be set for a single input file")

    # Extension → converter class name map (passed to workers)
    converters_map = {
        ".pdf": "PDFConverter",
        ".docx": "DocxConverter",
        ".pptx": "PptxConverter",
        ".ppt": "PptxConverter",
        ".doc": "DocxConverter",
    }

    # Build work items
    work_items = []
    for input_path in args.input:
        work_items.append((
            input_path,
            args.output if len(args.input) == 1 else None,
            args.image_dir,
            not args.no_images,
            args.verbose,
            args.force,
            converters_map,
        ))

    success = 0
    fail = 0
    tmp_files = []

    # Parallel processing (#11)
    if args.parallel and len(work_items) > 1:
        import multiprocessing
        workers = min(len(work_items), multiprocessing.cpu_count())
        if args.verbose:
            print(f"Parallel processing: {len(work_items)} files, {workers} workers")

        with multiprocessing.Pool(workers) as pool:
            results = pool.map(_convert_one_file, work_items)

        for ok, msg, tmp in results:
            print(msg)
            if ok:
                success += 1
            else:
                fail += 1
            if tmp:
                tmp_files.append(tmp)
    else:
        # Sequential processing
        for item in work_items:
            ok, msg, tmp = _convert_one_file(item)
            print(msg)
            if ok:
                success += 1
            else:
                fail += 1
            if tmp:
                tmp_files.append(tmp)

    # Cleanup temporary files
    import shutil as _shutil
    for tmp in tmp_files:
        try:
            _shutil.rmtree(os.path.dirname(tmp), ignore_errors=True)
        except Exception:
            pass

    # Summary
    if len(args.input) > 1:
        print(f"\nTotal: {success} succeeded, {fail} failed")

    sys.exit(0 if fail == 0 else 1)


if __name__ == "__main__":
    main()
